"""Build the JobHopper final-presentation deck.

Every slide is laid out against the grading rubric, and speaker notes carry the
narration plus the running clock.

Layout is measurement-driven: text is wrapped with real Carlito metrics (the
metric-compatible twin of Calibri) so a card is exactly as tall as the words
inside it and nothing silently overflows. build() prints a warning for any
slide whose content runs past the footer line.

Regenerate with:  python presentation/scripts/make_deck.py
"""
from __future__ import annotations

import os
from pathlib import Path

from PIL import ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
DIAGRAMS = ROOT / "diagrams"
OUT = ROOT / "JobHopper_Final_Presentation.pptx"

# ------------------------------------------------------------------ design --
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.72)
CONTENT_W = SLIDE_W - 2 * MARGIN
FOOTER_Y = SLIDE_H - Inches(0.62)

INK = RGBColor(0x1B, 0x24, 0x30)
MUTED = RGBColor(0x5E, 0x6B, 0x7C)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
CANVAS = RGBColor(0xF4, 0xF7, 0xFB)
HAIRLINE = RGBColor(0xC9, 0xD2, 0xDE)

BLUE = RGBColor(0x3F, 0x5C, 0x99)
TEAL = RGBColor(0x1F, 0x7A, 0x6B)
PLUM = RGBColor(0x9A, 0x4F, 0x70)
AMBER = RGBColor(0xB0, 0x74, 0x2C)
SLATE = RGBColor(0x5B, 0x64, 0x72)

FONT = "Calibri"
LINE = 1.22        # single-spaced line height, as a multiple of font size
EMU_PT = 12700

# Carlito, because it has the same glyph advances as Calibri: we lay the slides
# out by measuring Carlito and set the text in Calibri, and the two agree.
# Substituting any other font silently changes every measurement here, so a
# missing Carlito is an error rather than a fallback.
#
# Searched in order; set JOBHOPPER_DECK_FONT_DIR to a directory holding
# Carlito-Regular.ttf and Carlito-Bold.ttf to override. The path used to be a
# single Debian one, so nobody on a Mac could regenerate the deck at all.
_FONT_DIRS = [
    os.environ.get("JOBHOPPER_DECK_FONT_DIR", ""),
    "/usr/share/fonts/truetype/crosextra",                       # Debian/Ubuntu
    "/usr/share/fonts/carlito",                                  # Fedora/Arch
    "/Library/Fonts", "/Library/Fonts/Microsoft",                # macOS, system
    str(Path.home() / "Library/Fonts"),                          # macOS, user
    "/Applications/LibreOffice.app/Contents/Resources/fonts/truetype",
]
_FONT_FILES = {False: "Carlito-Regular.ttf", True: "Carlito-Bold.ttf"}
_CACHE: dict = {}
_WARN: list[str] = []


def _font_path(bold: bool) -> str:
    name = _FONT_FILES[bold]
    for directory in _FONT_DIRS:
        if directory and (Path(directory) / name).is_file():
            return str(Path(directory) / name)
    raise SystemExit(
        f"Cannot find {name}. The slide layout is measured in Carlito, which is\n"
        f"metric-compatible with Calibri, so it cannot be swapped for another font.\n"
        f"Install it (Debian/Ubuntu: 'apt install fonts-crosextra-carlito'; macOS:\n"
        f"download Carlito from Google Fonts into ~/Library/Fonts), or point\n"
        f"JOBHOPPER_DECK_FONT_DIR at a directory containing it.\n"
        f"Looked in: {', '.join(d for d in _FONT_DIRS if d)}"
    )


# ------------------------------------------------------------ measurement --
def _face(size_pt: float, bold: bool):
    key = (round(size_pt * 4), bold)
    if key not in _CACHE:
        _CACHE[key] = ImageFont.truetype(_font_path(bold), int(round(size_pt * 4)))
    return _CACHE[key]


def width_pt(text: str, size_pt: float, bold: bool = False) -> float:
    return _face(size_pt, bold).getlength(text) / 4.0


def wrap(text: str, avail_pt: float, size_pt: float, bold: bool = False
         ) -> list[str]:
    """Greedy wrap using real glyph advances. '\\n' forces a break."""
    out: list[str] = []
    for para in text.split("\n"):
        if not para.strip():
            out.append("")
            continue
        cur = ""
        for word in para.split():
            cand = f"{cur} {word}".strip()
            if cur and width_pt(cand, size_pt, bold) > avail_pt:
                out.append(cur)
                cur = word
            else:
                cur = cand
        if cur:
            out.append(cur)
    return out


def pt(length) -> float:
    return length / EMU_PT


# Renderers round differently than we measure, so a line measured at exactly the
# box width can wrap again and push a card over. Measure narrow, draw wide.
SLOP = Pt(7)


# ------------------------------------------------------------- primitives --
def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.25, rounded=False,
         radius=0.05):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        int(x), int(y), int(w), int(h),
    )
    shp.shadow.inherit = False
    # Drop the theme style reference; otherwise LibreOffice and Google Slides
    # re-apply the default shadow and outline even with an empty effect list.
    el = shp._element
    style = el.find("{http://schemas.openxmlformats.org/presentationml/2006/main}style")
    if style is not None:
        el.remove(style)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    if rounded:
        try:
            shp.adjustments[0] = radius
        except (IndexError, KeyError):
            pass
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return shp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def put(tf, text, size, *, bold=False, color=INK, first=False, space_before=0,
        space_after=0, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FONT
    return p


def canvas(slide):
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=CANVAS)


def footer(slide, number: str):
    tf = textbox(slide, MARGIN, FOOTER_Y + Inches(0.12), CONTENT_W, Inches(0.3))
    put(tf, "JobHopper   ·   The Laughing Bots", 11, color=MUTED, first=True)
    tf = textbox(slide, SLIDE_W - MARGIN - Inches(1.2), FOOTER_Y + Inches(0.12),
                 Inches(1.2), Inches(0.3))
    put(tf, number, 11, color=MUTED, first=True, align=PP_ALIGN.RIGHT)


def heading(slide, kicker, title, accent=BLUE, sub=None, title_size=34):
    y = Inches(0.5)
    tf = textbox(slide, MARGIN, y, CONTENT_W, Inches(0.24))
    put(tf, kicker.upper(), 13, bold=True, color=accent, first=True)

    lines = wrap(title, pt(CONTENT_W - SLOP), title_size, bold=True)
    th = Pt(len(lines) * title_size * LINE)
    tf = textbox(slide, MARGIN, y + Inches(0.28), CONTENT_W, th)
    for i, ln in enumerate(lines):
        put(tf, ln, title_size, bold=True, color=INK, first=(i == 0))

    bottom = y + Inches(0.28) + th + Pt(6)
    if sub:
        slines = wrap(sub, pt(CONTENT_W - SLOP), 17)
        sh = Pt(len(slines) * 17 * LINE)
        tf = textbox(slide, MARGIN, bottom, CONTENT_W, sh)
        for i, ln in enumerate(slines):
            put(tf, ln, 17, color=MUTED, first=(i == 0))
        bottom += sh + Pt(8)
    rect(slide, MARGIN, bottom, Inches(1.45), Pt(3.5), fill=accent)
    return bottom + Inches(0.3)


# ----------------------------------------------------------------- cards --
CARD_PAD_X = Inches(0.26)
CARD_PAD_T = Pt(13)
CARD_PAD_B = Pt(15)
CARD_GAP = Pt(8)
RAIL = Inches(0.055)


def card_metrics(w, title, body, title_size, body_size, tag=None):
    """Wrapped lines and the exact height a card needs."""
    inner = pt(w - RAIL - CARD_PAD_X - Inches(0.24) - SLOP)
    t_avail = inner - (58 if tag else 0)
    tl = wrap(title, t_avail, title_size, bold=True)
    bl = wrap(body, inner, body_size)
    h = (pt(CARD_PAD_T) + len(tl) * title_size * LINE + pt(CARD_GAP)
         + len(bl) * body_size * LINE + pt(CARD_PAD_B))
    return tl, bl, Pt(h)


def card(slide, x, y, w, h, accent, title, body, *, title_size=18,
         body_size=14, tag=None):
    tl, bl, _ = card_metrics(w, title, body, title_size, body_size, tag)
    rect(slide, x, y, w, h, fill=PAPER, line=HAIRLINE, rounded=True, radius=0.055)
    rect(slide, x, y, RAIL, h, fill=accent)

    tx = x + RAIL + CARD_PAD_X
    tw = w - RAIL - CARD_PAD_X - Inches(0.18)
    ty = y + CARD_PAD_T
    th = Pt(len(tl) * title_size * LINE)
    tf = textbox(slide, tx, ty, tw, th)
    for i, ln in enumerate(tl):
        put(tf, ln, title_size, bold=True, color=INK, first=(i == 0))

    by = ty + th + CARD_GAP
    tf = textbox(slide, tx, by, tw, Pt(len(bl) * body_size * LINE))
    blank_next = False
    for i, ln in enumerate(bl):
        if ln == "":
            blank_next = True
            continue
        put(tf, ln, body_size, color=MUTED, first=(i == 0),
            space_before=(body_size * 0.55 if blank_next else 0))
        blank_next = False

    if tag:
        badge = rect(slide, x + w - Inches(1.18), y + Pt(11), Inches(0.95),
                     Inches(0.25), fill=AMBER, rounded=True, radius=0.4)
        btf = badge.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        put(btf, tag, 9.5, bold=True, color=PAPER, first=True,
            align=PP_ALIGN.CENTER)


def card_grid(slide, y0, items, cols, *, title_size=19, body_size=14,
              gap_x=Inches(0.3), gap_y=Inches(0.24), uniform_rows=True):
    """Lay cards out in a grid, each row as tall as its tallest card."""
    w = (CONTENT_W - (cols - 1) * gap_x) / cols
    heights = [card_metrics(w, t, b, title_size, body_size, tag)[2]
               for _a, t, b, tag in items]
    y = y0
    for r in range(0, len(items), cols):
        row = items[r:r + cols]
        rh = max(heights[r:r + cols]) if uniform_rows else None
        for c, (accent, title, body, tag) in enumerate(row):
            h = rh if uniform_rows else heights[r + c]
            card(slide, MARGIN + c * (w + gap_x), y, w, h, accent, title, body,
                 title_size=title_size, body_size=body_size, tag=tag)
        y += (rh if uniform_rows else max(heights[r:r + cols])) + gap_y
    return y - gap_y


def banner(slide, y, text, size=17, fill=PAPER, color=INK, accent=None):
    lines = wrap(text, pt(CONTENT_W - Inches(0.72) - SLOP), size, bold=True)
    h = Pt(len(lines) * size * LINE) + Pt(26)
    rect(slide, MARGIN, y, CONTENT_W, h, fill=fill,
         line=(HAIRLINE if fill == PAPER else None), rounded=True, radius=0.14)
    if accent:
        rect(slide, MARGIN, y, RAIL, h, fill=accent)
    tf = textbox(slide, MARGIN + Inches(0.34), y + Pt(13),
                 CONTENT_W - Inches(0.68), Pt(len(lines) * size * LINE))
    for i, ln in enumerate(lines):
        put(tf, ln, size, bold=True, color=color, first=(i == 0))
    return y + h


def bullets(slide, x, y, w, items, size=20, gap=13):
    total = Pt(0)
    for i, item in enumerate(items):
        lines = wrap(item, pt(w) - 22 - pt(SLOP), size)
        h = Pt(len(lines) * size * LINE)
        rect(slide, x + Inches(0.02), y + Pt(size * 0.34), Pt(7), Pt(7),
             fill=TEAL, rounded=True, radius=0.5)
        tf = textbox(slide, x + Inches(0.26), y, w - Inches(0.26), h)
        for j, ln in enumerate(lines):
            put(tf, ln, size, color=INK, first=(j == 0))
        y += h + Pt(gap)
        total += h + Pt(gap)
    return y


def full_bleed(slide, image: Path, pad=Inches(0.16)):
    from PIL import Image

    with Image.open(image) as im:
        iw, ih = im.size
    aw, ah = SLIDE_W - 2 * pad, SLIDE_H - 2 * pad
    if iw / ih > aw / ah:
        w, h = aw, Emu(int(aw * ih / iw))
    else:
        h, w = ah, Emu(int(ah * iw / ih))
    slide.shapes.add_picture(str(image), Emu(int((SLIDE_W - w) / 2)),
                             Emu(int((SLIDE_H - h) / 2)), int(w), int(h))


def check(name: str, bottom):
    if bottom > FOOTER_Y:
        _WARN.append(f"{name}: content overruns footer by "
                     f"{pt(bottom - FOOTER_Y):.0f}pt")


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text.strip()


# ------------------------------------------------------------------ build --
def build() -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    # 1 -------------------------------------------------------------- title
    s = blank(prs)
    canvas(s)
    rect(s, 0, 0, Inches(0.32), SLIDE_H, fill=BLUE)
    rect(s, Inches(0.32), 0, Inches(0.1), SLIDE_H, fill=TEAL)
    rect(s, Inches(0.42), 0, Inches(0.06), SLIDE_H, fill=PLUM)

    tf = textbox(s, Inches(1.25), Inches(1.5), Inches(11.0), Inches(1.3))
    put(tf, "JobHopper", 66, bold=True, color=INK, first=True)
    tf = textbox(s, Inches(1.25), Inches(2.82), Inches(10.2), Inches(0.9))
    put(tf, "Turning a pile of job postings into a picture of what to learn "
            "— and a way to practise it.", 23, color=MUTED, first=True)

    rect(s, Inches(1.28), Inches(4.0), Inches(1.5), Pt(3.5), fill=TEAL)
    tf = textbox(s, Inches(1.25), Inches(4.36), Inches(11.0), Inches(0.3))
    put(tf, "THE LAUGHING BOTS", 14, bold=True, color=TEAL, first=True)
    tf = textbox(s, Inches(1.25), Inches(4.7), Inches(11.3), Inches(0.4))
    put(tf, "Rose Duarte   ·   Elijah Hewlett   ·   Angel Patino   ·   "
            "Terrell Whiting", 21, bold=True, color=INK, first=True)
    tf = textbox(s, Inches(1.25), Inches(5.32), Inches(10.6), Inches(0.9))
    put(tf, "Final project presentation", 17, color=MUTED, first=True)
    put(tf, "A full-stack web application: HTML, Bootstrap and Sass on the "
            "front end; FastAPI, SQLAlchemy and SQLite behind it.", 15,
        color=MUTED, space_before=7)
    footer(s, "1")
    notes(s, """
[0:00 - 0:15]   SPEAKER 1

Hi, we're the Laughing Bots - Rose, Elijah, Angel and Terrell. This is JobHopper:
it reads real job postings, shows you which skills those postings are actually
asking for, and lets you quiz yourself on any of them.
""")

    # 2 --------------------------------------------------------------- goal
    s = blank(prs)
    canvas(s)
    top = heading(s, "Application goal",
                  "The information is public. Reading it is the problem.", BLUE)
    col_w = (CONTENT_W - Inches(0.34)) / 2
    y = top
    for para in [
        "Somebody trying to move up in tech opens a job board and finds forty "
        "postings for the same title - each one a different wall of "
        "requirements.",
        "The useful signal is which tools keep coming up. That signal is real, "
        "it is public, and it is completely buried. Reading forty postings by "
        "hand is the thing nobody actually does, so people guess at what to "
        "learn next.",
    ]:
        lines = wrap(para, pt(col_w - SLOP), 19)
        h = Pt(len(lines) * 19 * LINE)
        tf = textbox(s, MARGIN, y, col_w, h)
        for i, ln in enumerate(lines):
            put(tf, ln, 19, color=INK, first=(i == 0))
        y += h + Pt(16)

    x2 = MARGIN + col_w + Inches(0.34)
    goal = ("JobHopper's goal is to analyze job postings for a target role, "
            "identify the most in-demand skills in the job market for that "
            "role, and present those findings to the user through a word cloud "
            "visual. Additionally, JobHopper aims to help people build those "
            "skills through interactive Q&A games.")
    _, _, gh = card_metrics(col_w, "What we set out to build", goal, 19, 15.5)
    card(s, x2, top, col_w, gh, TEAL, "What we set out to build", goal,
         title_size=19, body_size=15.5)

    b = banner(s, max(y, top + gh) + Inches(0.22),
               "And every skill in that picture is clickable - click it and "
               "you are in a timed quiz on that exact skill.", size=19,
               fill=INK, color=PAPER)
    check("s2", b)
    footer(s, "2")
    notes(s, """
[0:15 - 1:00]   SPEAKER 1

The problem isn't that this information is secret - job postings are public.
It's volume. Search one title and you get forty postings, each with its own wall
of requirements. The signal is which tools keep repeating, and nobody reads
forty postings to find it. So people guess at what to learn next.

So our goal was this: analyse the job postings for a target role, identify the
most in-demand skills, and show them as a word cloud - then help people actually
build those skills through interactive Q&A games.

That second half is what we care about. Every skill in the picture is clickable.
Click Python, and you're in a timed quiz on Python. It doesn't just tell you
what to learn - it gives you somewhere to start.
""")

    # 3 ------------------------------------------------ functional reqs 1/2
    s = blank(prs)
    canvas(s)
    top = heading(s, "Functional requirements  ·  1 of 2",
                  "Accounts, and building the word cloud", TEAL)
    cw2 = (CONTENT_W - Inches(0.3)) / 2
    fr1 = [
        (BLUE, "User registration and login",
         "The system shall allow users to register an account using a username "
         "and password input, with the option to include name and email.\n\n"
         "The system shall validate that the user's credentials are correct "
         "before allowing them to sign in.\n\n"
         "The system shall allow the user to see their most recent game keyword "
         "and score after the user logs in.\n\n"
         "The system shall allow the user to see their three most recent job "
         "title searches after the user logs in."),
        (TEAL, "Word cloud creation",
         "The system shall allow the user to input job search parameters "
         "including job title, location, and minimum salary.\n\n"
         "The system shall allow the user to select word cloud configurations "
         "such as shape and number of words.\n\n"
         "The system shall generate a word cloud where the size of each word is "
         "based on the frequency of tool and skill keywords extracted from "
         "online job postings."),
    ]
    h = max(card_metrics(cw2, t, b, 20, 14)[2] for _a, t, b in fr1)
    for i, (accent, title, body) in enumerate(fr1):
        card(s, MARGIN + i * (cw2 + Inches(0.3)), top, cw2, h, accent, title,
             body, title_size=20, body_size=14)
    check("s3", top + h)
    footer(s, "3")
    notes(s, """
[1:00 - 1:50]   SPEAKER 2

These are our functional requirements, written the way we specified them. I
won't read them all - two things to point at.

In registration and login, it's the last pair: showing you your most recent game
score, and your three most recent searches. Those are the requirements that
forced everything to be stored against a real account rather than sitting in the
browser.

And in word cloud creation, it's the last line. The size of each word is based
on how often that skill appears across the postings we pulled. Size equals
frequency - a count, not an opinion. That one sentence is the whole product.
""")

    # 4 ------------------------------------------------ functional reqs 2/2
    s = blank(prs)
    canvas(s)
    top = heading(s, "Functional requirements  ·  2 of 2",
                  "Scraping the postings, and the Q&A game", TEAL)
    scrape = ("The system shall calculate the frequency of tool and skill "
              "keywords extracted from online job postings.\n\n"
              "The system shall display an error in the event that there is not "
              "enough job information for the word cloud.")
    game = ("The system shall allow the user to click on any word in the word "
            "cloud.\n\n"
            "The system shall redirect the user to the respective Q&A game "
            "dashboard after the user clicks on a word from the word cloud.\n\n"
            "The system shall select Q&A game questions that are related to the "
            "skill or tool keyword and the difficulty that the user selected.\n\n"
            "The system shall keep a timer for the user's quiz completion based "
            "on the selected difficulty.\n\n"
            "The system shall display the score to the user.")
    _, _, h1 = card_metrics(CONTENT_W, "Job scraping", scrape, 20, 14)
    card(s, MARGIN, top, CONTENT_W, h1, PLUM, "Job scraping", scrape,
         title_size=20, body_size=14)
    y = top + h1 + Inches(0.24)
    _, _, h2 = card_metrics(CONTENT_W, "Q&A game", game, 20, 14)
    card(s, MARGIN, y, CONTENT_W, h2, PLUM, "Q&A game", game, title_size=20,
         body_size=14)
    check("s4", y + h2)
    footer(s, "4")
    notes(s, """
[1:50 - 2:40]   SPEAKER 2

Job scraping is two requirements. Count how often each keyword appears - and
display an error when there isn't enough job information to build a cloud. A
thin, misleading picture is worse than saying we can't build one, so we made
that a requirement rather than an edge case.

The Q&A game block is the exact path you'll see in the demo: click any word, get
sent to that word's game, get questions matching the keyword and difficulty you
chose, on a timer, with your score at the end.

One thing that isn't on the slide - every answer is locked in on the server
before you're told whether it was right, so the running score and the final
score can never disagree.
""")

    # 5 ------------------------------------------- non-functional reqs
    s = blank(prs)
    canvas(s)
    top = heading(s, "Non-functional requirements", "How well it has to do it",
                  PLUM,
                  sub="Not features - but the app is wrong without them.")
    end = card_grid(s, top, [
        (TEAL, "Word cloud speed",
         "The system should generate and display the skill word cloud within 10 "
         "seconds of the user submitting valid search criteria.", None),
        (TEAL, "Game responsiveness",
         "The system should redirect the user to the Q&A game within 10 "
         "seconds, and calculate their score within 7 seconds of finishing.",
         None),
        (BLUE, "Input safety",
         "The system should validate and sanitize all user input to prevent "
         "injection. Every query goes through the ORM, parameterised - never "
         "built by string.", None),
        (BLUE, "Credential security",
         "Passwords are bcrypt-hashed, never stored or returned. Tokens expire "
         "in 24 hours, and a failed login gives one generic message either way.",
         None),
        (PLUM, "Score integrity",
         "The server decides every score, never the browser, and a finished "
         "quiz cannot be replayed.", None),
        (SLATE, "Reliability",
         "250 automated tests - 222 backend, 28 front end - and nothing merges "
         "until both suites pass.", None),
    ], 3, title_size=19, body_size=13.5)
    check("s5", end)
    footer(s, "5")
    notes(s, """
[2:40 - 3:30]   SPEAKER 2

Non-functional requirements - not features, but the app is wrong without them.

The top row is what we committed to in our spec: the cloud back within ten
seconds, the game open within ten and scored within seven, and all user input
validated and sanitised against injection. We get that last one structurally -
every query goes through the ORM, parameterised, never built by string.

The bottom row we held ourselves to as we built. Passwords bcrypt-hashed and
never returned, and a failed login gives the same generic message whether the
account exists or not. The server decides every score, never the browser. And
250 tests, with nothing merging until both suites pass.
""")

    # 6 ----------------------------------------------------- stack diagram
    s = blank(prs)
    canvas(s)
    full_bleed(s, DIAGRAMS / "tech_stack.png")
    notes(s, """
[3:45 - 4:40]   SPEAKER 3

The whole system on one page.

Top: the browser. Twelve HTML pages, Bootstrap, Sass, plain JavaScript modules -
no framework, no build step. It talks to exactly one thing, our API, and
anything tied to an account carries a signed token.

The middle band is where every rule lives: FastAPI, seventeen endpoints, Pydantic
validating every request and response, and underneath it the services - hashing
and tokens, skill extraction, and the ingest that keeps our data current. Note
on the right that our postings are real, pulled from a live job-search API.

The bottom is one SQLite file, reached only through SQLAlchemy.

That line along the bottom is the rule that shaped everything: the browser never
touches the database, and never decides anything it could be lied to about.
""")
    footer(s, "6")

    # 7 ------------------------------------------------------ why the stack
    s = blank(prs)
    canvas(s)
    top = heading(s, "Why this stack",
                  "Chosen for a four-person team on a semester timeline", TEAL)
    end = card_grid(s, top, [
        (TEAL, "Python + FastAPI",
         "Python is the strongest shared language on this team. FastAPI gave us "
         "request and response validation for free and generated live API docs - "
         "so the front end could build against a written contract instead of "
         "waiting for the backend to finish.", None),
        (PLUM, "SQLite + SQLAlchemy",
         "SQLite is a single file: no database server to install, so four "
         "laptops and the CI runner all run an identical database with zero "
         "setup. SQLAlchemy keeps every query parameterised - injection-safe by "
         "construction - and moving to PostgreSQL later is one connection "
         "string.", None),
        (BLUE, "Figma, Bootstrap and Sass",
         "Figma let us agree on every screen before anyone implemented it, and "
         "kept us coordinated once we started. Bootstrap's components sped up "
         "the build and gave us a responsive grid on day one. Sass let us "
         "generate custom CSS quickly, from one shared palette across twelve "
         "pages.", None),
        (SLATE, "GitHub Actions + pytest + Jest",
         "The cheapest way to stop four people breaking each other's work. Both "
         "suites run on every pull request, so a regression is caught by the "
         "robot rather than by whoever pulls next.", None),
    ], 2, title_size=20)
    check("s7", end)
    footer(s, "7")
    notes(s, """
[4:40 - 5:35]   SPEAKER 3

Why these choices.

Python and FastAPI: Python is the strongest shared language on this team, and
FastAPI gave us validation for free plus live API docs - so the front end could
build against a written contract instead of waiting for the backend.

SQLite: one file, no server to install, so four laptops and the CI runner run an
identical database with zero setup - which removed a whole category of "works on
my machine". SQLAlchemy keeps every query parameterised, so we're injection-safe
by construction.

On the front end, Figma first. We designed every screen before anyone built it,
which kept four people from producing four different-looking pages. Bootstrap
gave us a responsive grid on day one, and Sass gave us one shared palette across
twelve pages.

And GitHub Actions - the cheapest way to stop four people breaking each other's
work.
""")

    # 8 ------------------------------------------------------ ERD overview
    s = blank(prs)
    canvas(s)
    full_bleed(s, DIAGRAMS / "erd_overview.png")
    notes(s, """
[5:50 - 6:30]   SPEAKER 4

Before the full diagram, the shape of it. Fourteen tables in three groups.

Left: the job-market data - roles, postings, skills, and the junctions between
them. That's where a cloud comes from. Middle: identity and activity. Right: the
quiz engine.

What makes this one system rather than three is two shared tables. Skills is
shared left to right - the same row that sizes a word in the cloud owns that
word's question bank, which is why clicking a word can start a quiz. And users
ties everything on the right to one account.
""")
    footer(s, "8")

    # 9 ---------------------------------------------------------- full ERD
    s = blank(prs)
    canvas(s)
    full_bleed(s, DIAGRAMS / "erd_full.png")
    notes(s, """
[6:15 - 7:35]   SPEAKER 4   -   Trace each table with the cursor as you name it.

Let me trace the one path that touches most of it.

Start at ROLES. One role has many JOB_POSTINGS - title, company, location, salary
range, date posted.

Now the interesting part. A posting mentions many skills, and a skill appears in
many postings. That's many-to-many, which a relational database can't store
directly - so JOB_SKILLS resolves it. Its primary key is the pair, job plus
skill, which is what guarantees one posting can only count once toward a skill.

So the word cloud is one query: the postings for this role inside the date
window, joined through job_skills, counted per skill, sorted descending. That
count is the word size.

Follow SKILLS right and it becomes the quiz. One skill has many QUESTIONS split
by difficulty, each with its ANSWER_OPTIONS - one flagged correct, and that flag
never leaves the server.

Starting a quiz creates a QUIZ_SESSION recording which questions went out and
what you picked - that is what makes live scoring trustworthy, and what stops a
finished quiz being replayed. Submitting lands the result in GAME_ATTEMPTS.

And down the middle, USERS - every search and every attempt hangs off it.
""")
    footer(s, "9")

    # 10 -------------------------------------------------------- the roles
    s = blank(prs)
    canvas(s)
    top = heading(s, "User roles",
                  "JobHopper has two roles — and we'll show you both", BLUE,
                  sub="What the app lets you do depends entirely on whether the "
                      "request carries a valid token.")
    cw2 = (CONTENT_W - Inches(0.3)) / 2
    pair = [
        (SLATE, "Visitor  ·  not signed in",
         "Can browse the home page, the game rules, the stack page and the "
         "creators page — and can register an account.\n\n"
         "Cannot generate a word cloud, and has no profile, no saved searches "
         "and no game history. The page sends them to sign-in — and the API "
         "refuses the request as well, so the rule holds even if you skip the "
         "page."),
        (BLUE, "Registered user  ·  signed in",
         "Everything above, plus: generate word clouds, play quizzes with the "
         "result saved, and a profile page showing recent searches and recent "
         "scores — each one re-runnable in a click.\n\n"
         "Identity comes from a signed token, checked on the server for every "
         "single request."),
    ]
    h = max(card_metrics(cw2, t, b, 21, 15)[2] for _a, t, b in pair)
    for i, (accent, title, body) in enumerate(pair):
        card(s, MARGIN + i * (cw2 + Inches(0.3)), top, cw2, h, accent, title,
             body, title_size=21, body_size=15)

    y = top + h + Inches(0.26)
    body = ("Nothing in JobHopper is administered through a screen. Postings "
            "arrive through the ingest pipeline and the question bank is loaded "
            "from a fixture, so there was never a task an admin would log in to "
            "do. Adding an admin login would have been a role with nothing "
            "behind it.")
    _, bl, ch = card_metrics(CONTENT_W, "Why there is no admin role", body, 19,
                             15)
    card(s, MARGIN, y, CONTENT_W, ch, AMBER, "Why there is no admin role", body,
         title_size=19, body_size=15)
    check("s10", y + ch)
    footer(s, "10")
    notes(s, """
[7:55 - 8:30]   SPEAKER 1   -   sets up the demo

One note before the demo, because it shapes what you'll see.

JobHopper has two roles, and the difference is whether the request carries a
valid token. A visitor can read the public pages and register - nothing else. A
signed-in user gets word clouds, quizzes that save their result, and a profile.

And there's deliberately no admin role. Postings come in through the ingest
pipeline and questions load from a fixture, so there was never a job an admin
would log in to do.
""")

    # 11-13 ------------------------------------------------------ demo cues
    demos = [
        ("1 OF 3", SLATE, "Role 1 — the visitor", [
            "Browse the public pages — home, rules, stack, creators",
            "Try to reach a word cloud without signing in → bounced to sign-in",
            "Register a brand-new account live — the form builds your first "
            "cloud with you",
        ], "11", """
[8:22 - 9:22]   SPEAKER 3   -   LIVE DEMO, PART 1.  Share your screen now.
Lines in quotes are what you say. Lines starting with > are what you do.

> Home page already open.

"This is JobHopper the way anyone arrives at it - not signed in. Rules, stack
and creators are open to everybody."

> Click Game Rules. Don't scroll far.

"Ten questions, three difficulties, each with its own clock."

> Try to open the word cloud page directly.

"Now watch what happens if I go straight for a word cloud with no account. It
sends me to sign-in - and that's not the page hiding a button. The API refuses
the request too, so there's no way around it."

> Click Sign Up. Type a 3-character username so the validation fires.

"So let's make one. Validation is inline, right where the problem is."

> Fix the username, then fill the rest. Job title and location are type-ahead
> fields - start typing and real options drop down. Keep talking while you
> type; don't narrate the typing.

"The same form takes your first search - and those suggestions are the job
titles and locations we actually have postings for."

> Submit.

If registration misbehaves: sign in with the backup account, say "we've already
got one set up", and carry on. Do not debug on camera.
"""),
        ("2 OF 3", BLUE, "Role 2 — the registered user", [
            "Read the cloud we just generated — then build another from the "
            "profile",
            "Click a skill → pick a difficulty → play the timed quiz",
            "Watch the live score move as each answer locks in",
            "Open the profile: the search and the score are both already there",
        ], "12", """
[9:22 - 11:22]   SPEAKER 3 and SPEAKER 4   -   LIVE DEMO, PART 2. The main event.
Lines in quotes are what you say. Lines starting with > are what you do.

> You are on the cloud registration just built.

"Here's what that search produced. The biggest words are the ones that appeared
in the most postings for that role - that's a count, not our opinion. Anything
we have questions for is clickable."

> Click one of the biggest skills.

"So let's take that one."

> Difficulty screen. Choose EASY - three minutes, so you can finish all ten.
> The score is only saved when the quiz is submitted, and that happens on the
> tenth answer or when the clock runs out. Stopping early means an empty game
> history on the profile.

"Three difficulties, three timers. I'll take easy - three minutes."

> Answer one question correctly.

"It tells me straight away whether I got it. It can do that safely because my
answer is already locked in on the server before I'm told - so knowing it now
can't change what I picked."

> Answer the next one wrong, on purpose.

"And there's the other case. It shows me the right answer, and the clock pauses
while I'm reading it."

> Now click through questions 3 to 10 without narrating each one. About
> sixty seconds if you don't read them aloud.

"I'll speed through the rest of these."

> The Submit button is replaced by a return button. Click it.

"That's graded and saved. This is my profile - the search I ran is under recent
word clouds, and the quiz I just played is under game history, score and all.
We didn't refresh anything, and none of this is hard-coded. That's the database
answering."

> Click Search Again on the saved cloud.

"And any saved search re-runs in one click."
"""),
        ("3 OF 3", PLUM, "It really is a database", [
            "Open the database file and list the tables — the same fourteen from "
            "the ERD",
            "Show the rows we just created: the search and the game attempt, "
            "carrying our user_id",
            "Show a stored password — a bcrypt hash, not the password we typed",
            "Show the word-cloud counts matching the picture we just drew",
        ], "13", """
[11:22 - 12:22]   SPEAKER 4   -   LIVE DEMO, PART 3. This is the 9-point item.
Switch to the DB Browser window already open on the second desktop. Do not open
a terminal - that reads as code.

> Show the table list in the sidebar.

"Last thing - proof this is really sitting in a database. There are our fourteen
tables, the same fourteen from the ERD. That diagram isn't what we planned; it's
what's in the file."

> DB Browser lists fifteen. The extra one is sqlite_stat1, which SQLite makes
> for its own query statistics. Say "our fourteen tables", not "fourteen
> tables", so the count on screen doesn't catch you out.

> Browse Data, searches table, scroll to the last row.

"The bottom row is the search from ninety seconds ago - role, salary, shape,
timestamp, and a user_id pointing at the account we made on camera."

> Switch to game_attempts.

"Same for the quiz. Skill, difficulty, score, seconds taken."

> Switch to users. Point at the password_hash column.

"And this is the users table. That's a bcrypt hash. We couldn't tell you that
password if we wanted to - we never stored it."

> Stop sharing, back to the deck.

"Everything you just saw came out of that one file."
"""),
    ]
    for part, accent, title, steps, num, note in demos:
        s = blank(prs)
        canvas(s)
        rect(s, 0, 0, SLIDE_W, Inches(0.13), fill=accent)
        tf = textbox(s, MARGIN, Inches(1.35), CONTENT_W, Inches(0.3))
        put(tf, f"LIVE DEMO  ·  PART {part}", 16, bold=True, color=accent,
            first=True)
        tlines = wrap(title, pt(CONTENT_W - SLOP), 50, bold=True)
        tf = textbox(s, MARGIN, Inches(1.78), CONTENT_W,
                     Pt(len(tlines) * 50 * LINE))
        for i, ln in enumerate(tlines):
            put(tf, ln, 50, bold=True, color=INK, first=(i == 0))
        ry = Inches(1.78) + Pt(len(tlines) * 50 * LINE) + Inches(0.18)
        rect(s, MARGIN, ry, Inches(1.5), Pt(3.5), fill=accent)
        end = bullets(s, MARGIN, ry + Inches(0.4), CONTENT_W - Inches(1.2),
                      steps, size=20)
        check(f"s{num}", end)
        footer(s, num)
        notes(s, note)

    # 14 --------------------------------------------- achieved vs altered
    s = blank(prs)
    canvas(s)
    top = heading(s, "Summary", "How much was achieved, and what was altered",
                  TEAL,
                  sub="We planned more than we shipped. Here is the honest "
                      "accounting.")
    cw2 = (CONTENT_W - Inches(0.3)) / 2
    pair = [
        (TEAL, "Achieved",
         "▸  Users can register, sign in, and stay signed in.\n\n"
         "▸  The app parses job postings and generates a word cloud from the "
         "user's own search parameters.\n\n"
         "▸  The Q&A game runs on a timer with live scoring, and every result "
         "is saved to the account.\n\n"
         "▸  The prep roadmap ranks a role's skills by demand and tracks "
         "progress through them.\n\n"
         "▸  The question bank beat its target: about 15 per skill per difficulty, "
         "which is 1,251 questions across 28 skills.\n\n"
         "▸  250 automated tests, run on every pull request."),
        (AMBER, "Altered or cut",
         "▸  Users were originally going to type in any role. External API "
         "limitations reduced that to 4 supported roles.\n\n"
         "▸  Location became a dynamic menu, listing only places that have "
         "enough postings to build a cloud from.\n\n"
         "▸  Maximum salary was removed rather than shipped as a filter that "
         "quietly did nothing.\n\n"
         "▸  Named ranks are on the rules page, but scoring currently ends at a "
         "normalised number."),
    ]
    h = max(card_metrics(cw2, t, b, 22, 14)[2] for _a, t, b in pair)
    for i, (accent, title, body) in enumerate(pair):
        card(s, MARGIN + i * (cw2 + Inches(0.3)), top, cw2, h, accent, title,
             body, title_size=22, body_size=14)
    check("s14", top + h)
    footer(s, "14")
    notes(s, """
[11:20 - 12:10]   SPEAKER 1

The honest accounting.

On the left, what worked. Everything you just saw, plus a question bank that beat
its own target - we aimed for ten per skill per difficulty and shipped fifteen,
so more than 1,250 questions. And 250 tests on every pull request.

On the right, what changed. The top two are the same story: the external API
limited what we could reliably fetch, so typing any role became four supported
roles, and location became a menu of places that actually have postings behind
them. Both were us refusing to offer a search that comes back empty.

The rest is choosing depth over breadth - and one place where our documentation
got ahead of the app, which is on us.
""")

    # 15 -------------------------------------------------------- close
    s = blank(prs)
    canvas(s)
    rect(s, 0, 0, Inches(0.32), SLIDE_H, fill=BLUE)
    rect(s, Inches(0.32), 0, Inches(0.1), SLIDE_H, fill=TEAL)
    rect(s, Inches(0.42), 0, Inches(0.06), SLIDE_H, fill=PLUM)

    left = Inches(1.25)
    inner_w = SLIDE_W - left - Inches(0.9)
    tf = textbox(s, left, Inches(0.85), inner_w, Inches(0.3))
    put(tf, "WHERE WE'D TAKE IT NEXT", 14, bold=True, color=TEAL, first=True)

    nx = [
        (BLUE, "Rank the roadmap by what you're weak at",
         "It orders skills by market demand today. It should also weigh what "
         "your quiz scores say you don't know yet."),
        (TEAL, "Widen the ingest",
         "Four roles was enough to prove the idea. The pipeline isn't the "
         "limit — the vocabulary of skills we match against is."),
        (PLUM, "Close the loop",
         "Feed quiz results back into the cloud, so the skills you keep getting "
         "wrong become the ones it puts in front of you."),
    ]
    cw3 = (inner_w - 2 * Inches(0.26)) / 3
    ch = max(card_metrics(cw3, t, b, 18, 13.5)[2] for _a, t, b in nx)
    for i, (accent, title, body) in enumerate(nx):
        card(s, left + i * (cw3 + Inches(0.26)), Inches(1.25), cw3, ch, accent,
             title, body, title_size=18, body_size=13.5)

    y = Inches(1.25) + ch + Inches(0.45)
    tf = textbox(s, left, y, inner_w, Inches(0.9))
    put(tf, "Thank you", 50, bold=True, color=INK, first=True)
    tf = textbox(s, left, y + Inches(0.86), inner_w - Inches(0.6), Inches(0.5))
    put(tf, "JobHopper reads what the market is asking for, and gives you "
            "somewhere to start.", 21, color=MUTED, first=True)
    rect(s, left + Inches(0.03), y + Inches(1.52), Inches(1.5), Pt(3.5),
         fill=TEAL)
    tf = textbox(s, left, y + Inches(1.85), inner_w, Inches(0.3))
    put(tf, "THE LAUGHING BOTS", 13, bold=True, color=TEAL, first=True)
    tf = textbox(s, left, y + Inches(2.18), inner_w, Inches(0.4))
    put(tf, "Rose Duarte   ·   Elijah Hewlett   ·   Angel Patino   ·   "
            "Terrell Whiting", 18, bold=True, color=INK, first=True)
    check("s15", y + Inches(2.5))
    footer(s, "15")
    notes(s, """
[12:30 - 13:00]   SPEAKER 1   -   close

Three things next. Rank the roadmap by what you're weak at, not just by what the
market wants. Widen the ingest; the pipeline isn't the limit, the skill
vocabulary is. And the one we actually want: close the loop, so the skills you
keep getting wrong become the ones your cloud puts in front of you.

Right now the two halves of the app share a database. They should share a memory.

That's JobHopper. Thank you.

--- END. Target 13:00. Hard ceiling 15:00. ---
""")

    prs.save(OUT)
    print(f"wrote {OUT}")
    if _WARN:
        print("LAYOUT WARNINGS:")
        for w in _WARN:
            print("  !", w)
    else:
        print("layout: all slides fit")


if __name__ == "__main__":
    build()
