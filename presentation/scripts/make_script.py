"""Extract the speaker notes from the built deck into SPEAKER_SCRIPT.md.

The deck is the single source of truth: the notes live in the .pptx, and this
regenerates the rehearsal document from them so the two can never drift.

It also estimates delivery time at 140 words per minute (an unhurried
presenting pace) and prints a warning if the total runs past the 15-minute
ceiling in the assignment.

Run after make_deck.py:  python presentation/scripts/make_script.py
"""
from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
DECK = ROOT / "JobHopper_Final_Presentation.pptx"
OUT = ROOT / "SPEAKER_SCRIPT.md"

WPM = 140
CEILING_S = 15 * 60
TARGET_S = 14 * 60

TITLES = [
    "Title", "The goal", "Functional requirements 1 of 2",
    "Functional requirements 2 of 2", "Non-functional requirements",
    "Technology stack diagram", "Why this stack", "ERD — subject areas",
    "ERD — full schema", "User roles", "Demo 1 — the visitor",
    "Demo 2 — the registered user", "Demo 3 — it really is a database",
    "Summary — achieved vs changed", "Close",
]

# Slides whose time is dominated by clicking, not talking.
DEMO_SLIDES = {11, 12, 13}


def spoken_words(note: str) -> int:
    """Count only the words that get said out loud."""
    words = 0
    for line in note.splitlines():
        st = line.strip()
        if not st:
            continue
        if st.startswith("["):                      # timing header
            continue
        if st.startswith("---"):                    # end marker
            continue
        if st.startswith("Runbook:"):
            continue
        if re.match(r"^\d+\.", st):                 # numbered demo steps
            continue
        if st.startswith(("If registration", "Close with:", "Do NOT", "Say:")):
            continue
        words += len(st.split())
    return words


def mmss(seconds: float) -> str:
    return f"{int(seconds // 60)}:{int(seconds % 60):02d}"


def main() -> None:
    prs = Presentation(DECK)
    rows = []
    running = 0.0
    body: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        note = ""
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
        title = TITLES[i - 1] if i <= len(TITLES) else f"Slide {i}"

        if i in DEMO_SLIDES:
            # Demo timing comes from the runbook, not from word count.
            est = {11: 90, 12: 135, 13: 70}[i]
        else:
            est = spoken_words(note) / WPM * 60
        running += est
        rows.append((i, title, est, running))

        header = re.search(r"^\[(.+?)\]\s*(.*)$", note, re.M)
        cue = header.group(2).strip(" —-") if header else ""
        clock = header.group(1) if header else ""

        body.append(f"## Slide {i} — {title}\n")
        if clock:
            body.append(f"**Clock:** `{clock}`  ·  **Speaker:** {cue or 'TBD'}  "
                        f"·  **Est. length:** {mmss(est)}\n")
        stripped = re.sub(r"^\[.*?\].*$", "", note, count=1, flags=re.M).strip()
        body.append(stripped + "\n")

    total = running
    lines = [
        "# JobHopper — final presentation speaker script",
        "",
        "> Generated from the speaker notes inside "
        "`JobHopper_Final_Presentation.pptx`.",
        "> Edit the notes in `scripts/make_deck.py`, rebuild the deck, then "
        "re-run `scripts/make_script.py`.",
        "",
        f"**Estimated total: {mmss(total)}**  "
        f"(target {mmss(TARGET_S)}, hard ceiling {mmss(CEILING_S)} — the "
        "assignment caps the video at 15 minutes; the rubric penalises going "
        "over 20.)",
        "",
        "Speaking estimate assumes 140 words per minute, which is an unhurried "
        "presenting pace. The three demo slides are timed from the runbook "
        "rather than from word count.",
        "",
        "## Running order",
        "",
        "| # | Slide | Length | Ends at |",
        "|---|-------|--------|---------|",
    ]
    for i, title, est, run in rows:
        lines.append(f"| {i} | {title} | {mmss(est)} | {mmss(run)} |")

    lines += [
        "",
        "## Who speaks",
        "",
        "The assignment does not require everyone to speak, but the rubric "
        "rewards a presentation that looks rehearsed. Suggested split — swap "
        "names to suit, and put whoever is most comfortable driving the app on "
        "the demo:",
        "",
        "| Speaker | Slides | Roughly |",
        "|---------|--------|---------|",
        "| Speaker 1 | 1, 2, 10, 14, 15 | open, goal, roles, summary, close |",
        "| Speaker 2 | 3, 4, 5 | functional and non-functional requirements |",
        "| Speaker 3 | 6, 7, 11, 12 | architecture, stack rationale, demo 1–2 |",
        "| Speaker 4 | 8, 9, 12, 13 | the database — ERD and the live data |",
        "",
        "## Delivery notes",
        "",
        "- **Look up.** Five rubric points ride on looking at the audience "
        "rather than reading. Learn the first and last sentence of each slide "
        "by heart; improvise the middle.",
        "- **No code on screen.** The assignment forbids showing code "
        "segments. The ERD, the architecture diagram and the database browser "
        "are all fine — an editor, a terminal or a source file is not.",
        "- **Say the numbers.** 14 tables, 1,260 questions, 40 postings, 237 "
        "tests. Specifics are what make a claim land.",
        "- **Do not narrate the mouse.** Say what a thing means, not where you "
        "are clicking.",
        "- **If something breaks on camera**, say what should have happened and "
        "move on. Do not debug live; stop the recording and re-take instead.",
        "",
        "---",
        "",
    ]
    lines += body

    OUT.write_text("\n".join(lines), encoding="utf-8")
    print(f"wrote {OUT}  —  estimated {mmss(total)}")
    if total > CEILING_S:
        print(f"  ! OVER the {mmss(CEILING_S)} ceiling — cut "
              f"{mmss(total - CEILING_S)}")
    elif total > TARGET_S:
        print(f"  ~ over target {mmss(TARGET_S)} but inside the ceiling")
    else:
        print(f"  ok — {mmss(CEILING_S - total)} of headroom under the ceiling")


if __name__ == "__main__":
    main()
